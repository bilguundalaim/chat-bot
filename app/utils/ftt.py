from google.cloud import storage
from google.oauth2 import service_account
from pptx import Presentation
from docx import Document
from app.utils.stt import transcribe_audio
import fitz
import os

def extract_text_from_pdf(file_path):
  text = ""
  with fitz.open(file_path) as doc:
    for page in doc:
      text += page.get_textbox()
  return text

def extract_text_from_pptx(file_path):
  prs = Presentation(file_path)
  text = ""
  for slide in prs.slides:
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        text += shape.text + "\n"
  return text

def extract_text_from_docx(file_path):
  doc = Document(file_path)
  text = ""
  for para in doc.paragraphs:
    text += para.text + "\n"
  return text

def download_blob_to_local(blob, download_path):
  os.makedirs(os.path.dirname(download_path), exist_ok=True)
  with open(download_path, "wb") as file:
    blob.download_to_file(file)

def download_and_extract_text(course_name):
  client_file = "./key.json"
  credentials = service_account.Credentials.from_service_account_file(client_file)
  client = storage.Client(credentials=credentials)

  bucket = client.bucket("lecture-files")
  blobs = bucket.list_blobs(prefix=f"{course_name}/")
  text = ""

  for blob in blobs:
    download_path = os.path.join("/tmp", os.path.basename(blob.name))

    download_blob_to_local(blob, download_path)

    # Extract text based on file type
    if blob.name.endswith('.pdf'):
      text += extract_text_from_pdf(download_path) + '\n'
    elif blob.name.endswith('.pptx'):
      text += extract_text_from_pptx(download_path) + '\n'
    elif blob.name.endswith('.docx'):
      text += extract_text_from_docx(download_path) + '\n'
    elif blob.name.endswith('.wav'):
      text += transcribe_audio(download_path) + '\n'
    else:
      print(f"Unsupported file type: {blob.name}")
    
  return text

def upload_to_bucket(source_file, destination_blob):
  client_file = "./key.json"
  credentials = service_account.Credentials.from_service_account_file(client_file)
  client = storage.Client(credentials=credentials)

  bucket = client.bucket("lecture-files")

  blob = bucket.blob(destination_blob)

  blob.upload_from_filename(source_file)

  blob.make_public()

  return blob.public_url